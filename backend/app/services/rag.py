import os
import hashlib
import asyncio
from typing import Optional
from sqlalchemy.ext.asyncio import AsyncSession
from app.config import settings


_file_cache: dict[str, list[dict]] = {}


class RAGService:
    def __init__(self, db: AsyncSession, user_id: Optional[str] = None):
        self.db = db
        self.user_id = user_id
        self.upload_dir = settings.upload_dir
        os.makedirs(self.upload_dir, exist_ok=True)
        self._qdrant = None
        self._embed_provider = None
        self._collection_ensured = False

    @staticmethod
    def cache_file(user_id: str, filename: str, file_path: str, text: str, file_id: str):
        if user_id not in _file_cache:
            _file_cache[user_id] = []
        _file_cache[user_id].append({
            "filename": filename,
            "file_path": file_path,
            "text": text,
            "file_id": file_id,
        })

    @staticmethod
    def get_cached_file_content(user_id: str, filename_hint: str) -> Optional[str]:
        files = _file_cache.get(user_id, [])
        for f in files:
            if filename_hint.lower() in f["filename"].lower():
                return f["text"]
        return None

    @staticmethod
    def get_cached_file_path(user_id: str, filename_hint: str) -> Optional[str]:
        files = _file_cache.get(user_id, [])
        for f in files:
            if filename_hint.lower() in f["filename"].lower():
                return f["file_path"]
        return None

    @staticmethod
    def get_all_cached_texts(user_id: str) -> Optional[str]:
        files = _file_cache.get(user_id, [])
        if not files:
            return None
        parts = []
        for f in files:
            parts.append(f"[From {f['filename']}]:\n{f['text']}")
        return "\n\n".join(parts)

    def _get_qdrant(self):
        if self._qdrant is None:
            try:
                from qdrant_client import AsyncQdrantClient
                self._qdrant = AsyncQdrantClient(url=settings.qdrant_url, timeout=5)
            except Exception:
                self._qdrant = False
        return self._qdrant if self._qdrant is not False else None

    def _get_embed_provider(self):
        if self._embed_provider is None:
            try:
                from app.services.nvidia.embeddings import NvidiaEmbeddingsProvider
                self._embed_provider = NvidiaEmbeddingsProvider()
            except Exception:
                self._embed_provider = False
        return self._embed_provider if self._embed_provider is not False else None

    async def _ensure_collection(self):
        if self._collection_ensured:
            return True
        qdrant = self._get_qdrant()
        if not qdrant:
            return False
        try:
            from qdrant_client.models import VectorParams, Distance
            collections = await asyncio.wait_for(
                qdrant.get_collections(), timeout=5.0
            )
            exists = any(c.name == settings.qdrant_collection for c in collections.collections)
            if not exists:
                await asyncio.wait_for(
                    qdrant.create_collection(
                        collection_name=settings.qdrant_collection,
                        vectors_config=VectorParams(size=4096, distance=Distance.COSINE),
                    ),
                    timeout=5.0,
                )
            self._collection_ensured = True
            return True
        except (asyncio.TimeoutError, Exception):
            return False

    async def process_file(self, file_path: str, filename: str, file_id: Optional[str] = None) -> dict:
        ext = os.path.splitext(filename)[1].lower()
        text = ""
        metadata = {"filename": filename, "path": file_path, "size": os.path.getsize(file_path)}

        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
        elif ext == ".md":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
        elif ext == ".pdf":
            text = self._extract_pdf(file_path)
        elif ext == ".docx":
            text = self._extract_docx(file_path)
        elif ext in (".csv", ".tsv"):
            text = self._extract_csv(file_path)
        elif ext in (".xlsx", ".xls"):
            text = self._extract_excel(file_path)
        elif ext == ".pptx":
            text = self._extract_pptx(file_path)
        elif ext in (".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".cs", ".go", ".rs", ".sql", ".html", ".css", ".json", ".xml", ".yaml", ".yml"):
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            metadata["language"] = ext.lstrip(".")
        else:
            text = f"File: {filename}\nType: {ext}\nSize: {metadata['size']} bytes\n"

        content_hash = hashlib.sha256(text.encode()).hexdigest()
        chunks = self._chunk_text(text)

        if file_id and text.strip():
            await self._index_chunks(chunks, file_id, filename, metadata)

        return {
            "text": text,
            "chunks": chunks,
            "metadata": metadata,
            "content_hash": content_hash,
            "chunk_count": len(chunks),
        }

    async def _index_chunks(self, chunks: list[str], file_id: str, filename: str, metadata: dict):
        embed_provider = self._get_embed_provider()
        qdrant = self._get_qdrant()
        if not embed_provider or not qdrant:
            return
        if not await self._ensure_collection():
            return

        from qdrant_client.models import PointStruct

        batch_size = 10
        points = []
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i + batch_size]
            try:
                embeddings = await embed_provider.create(texts=batch, input_type="passage")
            except Exception:
                continue
            for j, (chunk_text, embedding) in enumerate(zip(batch, embeddings)):
                points.append(PointStruct(
                    id=f"{file_id}_{i + j}",
                    vector=embedding,
                    payload={
                        "file_id": file_id,
                        "filename": filename,
                        "chunk_index": i + j,
                        "text": chunk_text,
                        "user_id": self.user_id or "",
                    },
                ))
        if points:
            try:
                await asyncio.wait_for(
                    qdrant.upsert(
                        collection_name=settings.qdrant_collection,
                        points=points,
                    ),
                    timeout=5.0,
                )
            except (asyncio.TimeoutError, Exception):
                pass

    async def search_similar(self, query: str, top_k: int = 5) -> Optional[str]:
        if not query or not query.strip():
            return None
        embed_provider = self._get_embed_provider()
        qdrant = self._get_qdrant()
        if not embed_provider or not qdrant:
            return None
        try:
            query_embedding = await embed_provider.create(texts=[query], input_type="query")
        except Exception:
            return None
        if not query_embedding:
            return None

        from qdrant_client.models import Filter, FieldCondition, MatchValue

        filter_ = None
        if self.user_id:
            filter_ = Filter(
                must=[FieldCondition(key="user_id", match=MatchValue(value=self.user_id))]
            )

        try:
            results = await asyncio.wait_for(
                qdrant.search(
                    collection_name=settings.qdrant_collection,
                    query_vector=query_embedding[0],
                    limit=top_k,
                    query_filter=filter_,
                    score_threshold=0.5,
                ),
                timeout=5.0,
            )
        except (asyncio.TimeoutError, Exception):
            return None

        if not results:
            return None

        context_parts = []
        for r in results:
            text = r.payload.get("text", "")
            filename = r.payload.get("filename", "")
            if text:
                context_parts.append(f"[From {filename}]: {text}")

        if not context_parts:
            return None

        return "\n\n".join(context_parts)

    def _chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> list[str]:
        if len(text) <= chunk_size:
            return [text]
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            if end < len(text):
                last_period = text.rfind(".", start, end)
                last_newline = text.rfind("\n", start, end)
                split_at = max(last_period, last_newline)
                if split_at > start:
                    end = split_at + 1
            chunks.append(text[start:end])
            start = end - overlap if end < len(text) else len(text)
        return chunks

    def _extract_pdf(self, path: str) -> str:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            return f"[PDF file: {os.path.basename(path)}]"

    def _extract_docx(self, path: str) -> str:
        try:
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return f"[DOCX file: {os.path.basename(path)}]"

    def _extract_csv(self, path: str) -> str:
        import pandas as pd
        df = pd.read_csv(path)
        return df.to_string()

    def _extract_excel(self, path: str) -> str:
        import pandas as pd
        dfs = pd.read_excel(path, sheet_name=None)
        parts = []
        for sheet_name, df in dfs.items():
            parts.append(f"--- Sheet: {sheet_name} ---\n{df.to_string()}")
        return "\n\n".join(parts)

    def _extract_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception:
            return f"[PPTX file: {os.path.basename(path)}]"
